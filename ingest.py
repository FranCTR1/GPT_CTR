import os
from dotenv import load_dotenv
from openai import OpenAI
import chromadb
from PyPDF2 import PdfReader
from pptx import Presentation
import pandas as pd
import docx

# Cargar variables de entorno
load_dotenv()
api_key = os.getenv("OPENAI_API_KEY")

# Inicializar clientes
client_openai = OpenAI(api_key=api_key)
client_chroma = chromadb.PersistentClient(path="./chromadb_local")
collection = client_chroma.get_or_create_collection("empresa_docs")

# Funciones para extracción de texto
def extract_text_from_pdf(path):
    reader = PdfReader(path)
    return [page.extract_text() for page in reader.pages if page.extract_text()]

def extract_text_from_csv(path):
    df = pd.read_csv(path)
    rows = []
    for idx, row in df.iterrows():
        formatted = "\n".join([f"{col.strip()}: {str(val).strip()}" for col, val in row.items()])
        rows.append(formatted)
    return rows

def extract_text_from_excel(path):
    dfs = pd.read_excel(path, sheet_name=None)
    rows = []
    for sheet, df in dfs.items():
        for _, row in df.iterrows():
            formatted = "\n".join([f"{col.strip()}: {str(val).strip()}" for col, val in row.items()])
            rows.append(formatted)
    return rows

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
    return text

def extract_text_from_docx(path):
    doc = docx.Document(path)
    return [para.text.strip() for para in doc.paragraphs if para.text.strip()]

def extract_text_from_document(file_path):
    ext = file_path.lower().split('.')[-1]
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext == "pptx":
        return extract_text_from_pptx(file_path)
    elif ext == "xlsx":
        return extract_text_from_excel(file_path)
    elif ext == "docx":
        return extract_text_from_docx(file_path)
    elif ext == "csv":
        return extract_text_from_csv(file_path)
    else:
        raise ValueError(f"Formato no soportado: {ext}")

# Recorrer documentos en la carpeta /data
documentos_dir = "data"
documentos = [os.path.join(documentos_dir, file) for file in os.listdir(documentos_dir)]

for doc in documentos:
    try:
        print(f"📁 Procesando: {os.path.basename(doc)}")
        fragments = extract_text_from_document(doc)

        for idx, fragment in enumerate(fragments):
            if not fragment.strip():
                continue

            # Generar embedding
            response = client_openai.embeddings.create(
                input=fragment,
                model="text-embedding-ada-002"
            )
            embedding = response.data[0].embedding

            # Indexar en ChromaDB
            collection.add(
                documents=[fragment],
                embeddings=[embedding],
                ids=[f"{os.path.basename(doc)}_{idx}"]
            )

        print(f"✅ {os.path.basename(doc)} cargado con {len(fragments)} fragmentos.")

    except Exception as e:
        print(f"⚠️ Error al procesar {doc}: {e}")

print("🎉 Ingesta finalizada: Todos los documentos fueron indexados.")
